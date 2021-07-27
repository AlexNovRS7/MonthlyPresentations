from pptx import Presentation
import calendar
import os 
import shutil

from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION

from pptx.dml.color import RGBColor
from pptx.util import Pt

from countries import countries, units

import govcharts

import datetime
import time

import pyodbc

today = datetime.date.today()
first = today.replace(day=1)
sqldate = f'{first.year}-{first.month - 1}-{first.day}' # date for SQL querries
lastMonth = first - datetime.timedelta(days=1)


# conn = pyodbc.connect('Driver={SQL Server};'
#                       'Server=XE-S-CWPDB01P.XE.ABB.COM;'
#                       'Database=PG_AskHR_KPI_2;'
#                       'Trusted_Connection=yes;')

# cursor = conn.cursor()

def CreateReport(ctry):
    t1 = time.perf_counter()

    conn = pyodbc.connect('Driver={SQL Server};'
                      'Server=XE-S-CWPDB01P.XE.ABB.COM;'
                      'Database=PG_AskHR_KPI_2;'
                      'Trusted_Connection=yes;')

    cursor = conn.cursor()

    global countries
    print("Creating: ", countries[ctry])

    prs = Presentation(f"{os.getcwd()}\\Created\\GovDecks\\{countries[ctry]}.pptx")
    slides = prs.slides
    
    # slide 1 (Cover)
    slide = slides[0]
    shapes = slide.shapes
    for shape in shapes:
        if shape.name == 'ctry_range_title':
            shape.text = 'Performance Management Report – ' +  countries[ctry] + '\nHR Operations Quality and Continuous Improvement'
        
        if shape.name == 'date_range':
            text = str(today)
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            
            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)
                        
    
    
    # slide 2 (Exec Summary)
    shapes = slides[1].shapes
    for shape in shapes:
            if shape.name == 'ctry_range':
                text = f'Performance {countries[ctry]} at a glance – HR Center and FO view'
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
            
                font = run.font
                font.name = 'Arial'
                font.size = Pt(14)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
            # add CSAT chart
            elif shape.name == 'csatChart_area':
                months, CSAT, ResponseRate = govcharts.get_cs_chart(cursor, sqldate, ctry, "%", "%")
            
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('CSAT', CSAT)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[1]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0"%"'

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)

            # add OTD chart
            elif shape.name == 'otdChart_area':
                months, OTD = govcharts.get_otd_chart(cursor, sqldate, ctry, "%")
            
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('OTD', OTD)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[1]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0"%"'

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
            
            # add Quality chart
            elif shape.name == 'qcChart_area':
                months, QC = govcharts.get_qc_chart(cursor, sqldate, ctry, "%")
            
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('QC', QC)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[1]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0.0"%"'

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
           
            # add Case Ageing chart
            elif shape.name == 'caChart_area':
                months, CA = govcharts.get_ca_chart(cursor, sqldate, ctry, "%")
            
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('CA', CA)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[1]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0"%"'

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
            
    # slides 3, 14 (Section name)
    for s in (2, 13):
        if s == 2:
            unit = 'Hub_Office'
        elif s == 13:
            unit = 'Front_Office'
         
        shapes = slides[s].shapes
        for shape in shapes:
            if shape.name == 'ctry_range':
                text = f"HR Operations {countries[ctry]} Monthly Review | {units[unit]}"
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
        
                font = run.font
                font.name = 'Arial'
                font.size = Pt(48)
                font.bold = True
                font.color.rgb = RGBColor(0, 0, 0)
                
                
    # slide 4, 14  (Summary)
    for s in (3, 13):
        if s == 3:
            unit = 'Hub_Office'
        elif s == 13:
            unit = 'Front_Office'
        feedbackreceived, rating12Nom, createdC, createdU, resolvedU, \
        openworkload, ageing115, ageing1530, ageing3060, ageing60180, \
        ageingover180, otdNom, otdDenom, QCDoneNom, QCnaNom, QCpassedNom, \
        QCfailedNom, QCDoneNomVol, QCnaNomVol, QCpassedNomVol, QCfailedNomVol = govcharts.get_summary_data(cursor, sqldate, ctry, unit)
        
        shapes = slides[s].shapes
        for shape in shapes:
            if shape.name == 'month_range':
                text = calendar.month_name[lastMonth.month] + ' ' + str(lastMonth.year)
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
        
                font = run.font
                font.name = 'Arial'
                font.size = Pt(18)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
        
            if shape.name == 'ctry_range':
                    text = f'Performance {countries[ctry]} at a glance – {units[unit]} view'
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = text
            
                    font = run.font
                    font.name = 'Arial'
                    font.size = Pt(14)
                    font.bold = False
                    font.color.rgb = RGBColor(0, 0, 0)
    
            elif shape.name == 'CFArea':
                satisfaction = round((rating12Nom / feedbackreceived) * 100) if feedbackreceived > 0 else 100
                response_rate = round((feedbackreceived * 100) / createdC) if createdC > 0 else 100

                dissatisfaction = 100 - satisfaction
                text = f"""{satisfaction}% ({rating12Nom}) Very good / Good;
{dissatisfaction}% ({feedbackreceived - rating12Nom}) not satisfied
Response rate: {response_rate}% ({feedbackreceived} answers)"""
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
        
                font = run.font
                font.name = 'Arial'
                font.size = Pt(14)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
    
            elif shape.name == 'VolumeArea':
                text = f"""{createdU} created / {resolvedU} closed
       """
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
        
                font = run.font
                font.name = 'Arial'
                font.size = Pt(14)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
    
            elif shape.name == 'CAArea':
                u15d = round(ageing115 * 100 / openworkload) if openworkload > 0 else 0
                u30d = round(ageing1530 * 100 / openworkload) if openworkload > 0 else 0
                u60d = round(ageing3060 * 100 / openworkload) if openworkload > 0 else 0
                u180d = round(ageing60180 * 100 / openworkload) if openworkload > 0 else 0
                a180d = round(ageingover180 * 100 / openworkload) if openworkload > 0 else 0

                text = f"""Open workload {openworkload}
{u15d}% ageing <15d / {u30d}% 15-30d
/ {u60d}% 30-60d / {u180d}% 60-180d / {a180d}% >180d"""
                
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
        
                font = run.font
                font.name = 'Arial'
                font.size = Pt(14)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
    
            elif shape.name == 'OTDArea':
                _otd = round(otdNom * 100 / otdDenom) if otdDenom > 0 else 100

                text = f"""{_otd}% On Time Delivery
"""
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
        
                font = run.font
                font.name = 'Arial'
                font.size = Pt(14)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
        
            elif shape.name == 'QCAreaTOP':
                text = f"""Out of all survey responses ({feedbackreceived})
X% (Y) was unsatisfied with Quality"""
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
        
                font = run.font
                font.name = 'Arial'
                font.size = Pt(14)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
        
            elif shape.name == 'QCArea':
                done = round(QCDoneNom * 100 / (QCDoneNom + QCnaNom)) if QCDoneNom + QCnaNom > 0 else 0
                passed = round(QCpassedNom * 100 / QCDoneNom) if QCDoneNom > 0 else 0
                failed = round(QCfailedNom * 100 / (QCDoneNom + QCnaNom)) if QCDoneNom + QCnaNom > 0 else 0
                not_applies = round(QCnaNom * 100 / (QCDoneNom + QCnaNom)) if QCDoneNom + QCnaNom > 0 else 0

                text = f"""{done}% ({QCDoneNomVol}) QC done, out of it:
{passed}% ({QCpassedNomVol}) passed / {failed}% ({QCfailedNomVol}) failed
{not_applies} % ({QCnaNomVol}) QC n/a"""
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
        
                font = run.font
                font.name = 'Arial'
                font.size = Pt(14)
                font.bold = False
                font.color.rgb = RGBColor(0, 0, 0)
    
    # Service Lines slides
    SLs = ["TA", "ELCM", "PY", "CB", "LD", "TM", "PA", "T1"]
    slides_ = [int(x) for x in range(4,11)]
    slides_.extend([int(x) for x in range(14,22)])
    for s in slides_:
        shapes = slides[s].shapes
        shapes = [shape for shape in shapes if "Area" in shape.name or "range" in shape.name]
        if s < 12:
            SL = SLs[s - 4]
            unit = 'Hub_Office'
        else:
            SL = SLs[s - 14]
            unit = 'Front_Office'
        for shape in shapes:
            if shape.name == 'ctry_range':
                curr_text = str(shape.text)
                a,b = curr_text.split(' Country ')
                new_text = a + f' {countries[ctry]} ' + b
                shape.text = new_text
            
            if shape.name == 'month_range':
                shape.text = 'Performance Review for ' + calendar.month_name[lastMonth.month] + ' ' + str(lastMonth.year)
            
            if shape.name == 'CSChart_Area':
                months, CSAT, ResponseRate = govcharts.get_cs_chart(cursor, sqldate, ctry, SL, unit)
            
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('CSAT', CSAT)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[s]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0"%"'

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
            # OTD
            elif shape.name == 'OTDChart_Area':
                months, OTD = govcharts.get_otd_chart(cursor, sqldate, ctry, SL, unit)
            
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('OTD', OTD)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[s]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0"%"'

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
            
            # add Quality chart
            elif shape.name == 'QCChart_Area':
                months, QC = govcharts.get_qc_chart(cursor, sqldate, ctry, SL, unit)
            
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('QC', QC)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[s]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                plot = chart.plots[0]
                plot.has_data_labels = True
                plot.gap_width = 100
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0.0"%"'

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)

            # add Volumes chart
            elif shape.name == 'VLChart_Area':
                months, created, resolved = govcharts.get_vl_chart(cursor, sqldate, ctry, SL, unit)

                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('Created', created)
                chart_data.add_series('Resolved', resolved)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[s]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(8)
                #chart.chart_style = 1

                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(8)
                data_labels.font.color.rgb = RGBColor(255, 255, 255)
                data_labels.position = XL_LABEL_POSITION.CENTER

                value_axis = chart.value_axis
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)
                value_axis.has_major_gridlines = False
                chart.value_axis.visible = False

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)

            # add Case ageing chart
            elif shape.name == 'CAChart_Area':
                months, active115, active1530, active3060, activeover60, longrunning = govcharts.get_ca_chart(cursor, sqldate, ctry, SL, unit)
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
            
                for c in chart_data.categories:
                    c.add_sub_category('active')
                    c.add_sub_category('on hold')
                

                
                chart_data.add_series('Long running cases', longrunning)
                chart_data.add_series('1 - 15 days', active115)
                chart_data.add_series('15 - 30 days', active1530)
                chart_data.add_series('30 - 60 days', active3060)
                chart_data.add_series('over 60 days', activeover60)
                
            
                
                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[s]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
                )
                
                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_title = False
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(8)
                # chart.chart_style = 1
                

                plot = chart.plots[0]
                plot.has_data_labels = True
                plot.gap_width = 75
                data_labels = plot.data_labels
                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(255, 255, 255)
                data_labels.position = XL_LABEL_POSITION.CENTER

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(8)
   
    # slide 13, 25 (Exec Summary)
    for s in (12, 22, 23, 24):
        shapes = slides[s].shapes
        for shape in shapes:
                if shape.name == 'ctry_range':
                    curr_text = str(shape.text)
                    a,b = curr_text.split(' Country ')
                    new_text = a + f' {countries[ctry]} ' + b
                    shape.text = new_text

    # save presentation
    name = 'HR Operations Monthly Review ' + countries[ctry] + ' ' + calendar.month_name[lastMonth.month] + ' ' + str(
        lastMonth.year) + '.pptx'
        
    # prs.save(f"{os.getcwd()}\\Created\\GovDecks\\" + name)
    prs.save(f"{os.getcwd()}\\Created\\GovDecks\\" + name)
    os.remove(f"{os.getcwd()}\\Created\\GovDecks\\{countries[ctry]}.pptx")
        
    t2 = time.perf_counter()
    print(countries[ctry], f' created. \t time: {round((t2 - t1), 2)} seconds')

    cursor.close()

