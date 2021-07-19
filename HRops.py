from pptx import Presentation
import calendar

import os
import pyodbc
import datetime
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION

from pptx.dml.color import RGBColor
from pptx.util import Pt

from pptx.chart.data import CategoryChartData

import time


def get_summary_data(cursor, sqldate, unit):
    arguments = f"'{sqldate}','{unit}'"
    cursor.execute('exec [HRmr_get_summarySlide]' + arguments)
    for row in cursor:
        feedbackreceived = row[1]
        rating12Nom = row[0]
        createdC = row[2]
        createdU = row[3]
        resolvedC = row[4]
        resolvedU = row[5]
        openworkload = row[6]
        ageing115 = row[7]
        ageing1530 = row[8]
        ageing3060 = row[9]
        ageing60180 = row[11]
        ageingover180 = row[12]
        otdNom = row[13]
        otdDenom = row[14]
        QCDoneNom = row[15]
        QCnaNom = row[16]
        QCpassedNom = row[17]
        QCfailedNom = row[18]
        QCDoneNomVol = row[19]
        QCnaNomVol = row[20]
        QCpassedNomVol = row[21]
        QCfailedNomVol = row[22]
    return feedbackreceived, rating12Nom, createdC, createdU, resolvedC,resolvedU, \
           openworkload, ageing115, ageing1530, ageing3060, ageing60180, \
           ageingover180, otdNom, otdDenom, QCDoneNom, QCnaNom, QCpassedNom, \
           QCfailedNom, QCDoneNomVol, QCnaNomVol, QCpassedNomVol, QCfailedNomVol


def get_query_info(name):
    if name in ('BGL', 'KRK', 'SPL', 'TLL'):
        typeN = 'pHub'
    elif name in ('CSE', 'NEU', 'NAM', 'SAM', 'NEA', 'MEA', 'SEA'):
        typeN = 'pRegion'
    elif name in ('TA_', 'TM_', 'GM_', 'PA_', 'ELC', 'CB_', 'PY_', 'LD_'):
        typeN = 'pSline'
        if name == 'ELC':
            name = 'ELCM'
        else:
            name = name[0:2]
    return name, typeN


def get_cs_chart(cursor, sqldate, name):
    name, typeN = get_query_info(name)
    cursor.execute(f"exec [HRmr_get_CS_chart] '{sqldate}','{name}','{typeN}'")

    months = []
    VeryGood = []
    Good = []
    Dissatisfied = []
    HighlyDissatisfied = []
    ResponseRate = []

    for row in cursor:
        months.append(row[0])
        VeryGood.append(row[1])
        Good.append(row[2])
        Dissatisfied.append(row[3])
        HighlyDissatisfied.append(row[4])
        if row[6] == 0:
            ResponseRate.append(0)
        else:
            ResponseRate.append(round((row[5]/row[6] * 100) ,1))
    
    VeryGood = [x if x > 0 else "" for x in VeryGood]
    Good = [x if x > 0 else "" for x in Good]
    Dissatisfied = [x if x > 0 else "" for x in Dissatisfied]
    HighlyDissatisfied = [x if x > 0 else "" for x in HighlyDissatisfied]
    return months, VeryGood, Good, Dissatisfied, HighlyDissatisfied, ResponseRate


def get_otd_chart(cursor, sqldate, name):
    name, typeN = get_query_info(name)
    cursor.execute(f"exec [HRmr_get_OTD_chart] '{sqldate}','{name}','{typeN}'")

    months = []
    data = []

    for row in cursor:
        months.append(row[0])
        if row[2] == 0:
            data.append(0)
        else:
            data.append(round(row[1] / row[2] * 100, 2))
    
    
    return months, data


def get_qc_chart(cursor, sqldate, name):
    name, typeN = get_query_info(name)
    cursor.execute(f"exec [HRmr_get_QC_chart] '{sqldate}','{name}','{typeN}'")

    months = []
    data = []

    for row in cursor:
        months.append(row[0])
        if row[3] == 0:
            data.append(0)
        else:
            data.append(row[1] * 100 / row[3])
    return months, data


def get_vl_chart(cursor, sqldate, name):
    name, typeN = get_query_info(name)
    cursor.execute(f"exec [HRmr_get_VL_chart] '{sqldate}','{name}','{typeN}'")

    months = []
    created = []
    resolved = []

    for row in cursor:
        months.append(row[0])
        created.append(row[1])
        resolved.append(row[2])
    return months, created, resolved


def get_ca_chart(cursor, sqldate, name):
    name, typeN = get_query_info(name)
    cursor.execute(f"exec [HRmr_get_CA_chart] '{sqldate}','{name}','{typeN}'")
    
    months = []
    active115 = []
    active1530 = []
    active3060 = []
    activeover60 = []
    onhold115 = []
    onhold1530 = []
    onhold3060 = []
    onhold60 = []
    longrunning = []
    
    for row in cursor:
        months.append(row[0])
        active115.append(row[1])
        active1530.append(row[2])
        active3060.append(row[3])
        activeover60.append(row[4])
        onhold115.append(row[5])
        onhold1530.append(row[6])
        onhold3060.append(row[7])
        onhold60.append(row[8])
        longrunning.append(row[9])
    
    t115 = [None]*(len(active115)+len(onhold115))
    t115[::2] = active115
    t115[1::2] = onhold115
    
    t1530 = [None]*(len(active115)+len(onhold115))
    t1530[::2] = active1530
    t1530[1::2] = onhold1530
    
    t3060 = [None]*(len(active115)+len(onhold115))
    t3060[::2] = active3060
    t3060[1::2] = onhold3060
    
    to60 = [None]*(len(active115)+len(onhold115))
    to60[::2] = activeover60
    to60[1::2] = onhold60
    
    lr = [0] * (len(longrunning) * 2)
    lr[1::2] = longrunning
    
    t115 = [x if x > 0 else "" for x in t115]
    t1530 = [x if x > 0 else "" for x in t1530]
    t3060 = [x if x > 0 else "" for x in t3060]
    to60 = [x if x > 0 else "" for x in to60]
    lr = [x if x > 0 else "" for x in lr]
    
    return months, t115, t1530, t3060, to60, lr 


def get_calt_chart(cursor, sqldate, name):
    name, typeN = get_query_info(name)
    cursor.execute(f"exec [HRmr_get_CALT_chart] '{sqldate}','{name}','{typeN}'")
    
    months = []
    ag120150 = []
    ag150180 = []
    agover180 = []
    
    for row in cursor:
        months.append(row[0])
        ag120150.append(row[1]* 100)
        ag150180.append(row[2] * 100)
        agover180.append(row[3] * 100)
    
    ag120150 = [x if x > 0 else "" for x in ag120150]
    ag150180 = [x if x > 0 else "" for x in ag150180]
    agover180 = [x if x > 0 else "" for x in agover180]  
    
    return months, ag120150, ag150180, agover180

def CreateReport():
    start_time = time.time()

    # open file
    prs = Presentation(f"{os.getcwd()}\\templates\\Monthly_HR_Report_TEMPLATE_4.0.pptx")

    today = datetime.date.today()
    first = today.replace(day=1)
    lastMonth = first - datetime.timedelta(days=1)
    sqldate = f'{first.year}-{first.month - 1}-{first.day}' # date for SQL querries
    slides = prs.slides

    conn = pyodbc.connect('Driver={SQL Server};'
                        'Server=XE-S-CWPDB01P.XE.ABB.COM;'
                        'Database=PG_AskHR_KPI_2;'
                        'Trusted_Connection=yes;')

    cursor = conn.cursor()

    # slide 1 (Cover)
    slide = slides[0]
    shapes = slide.shapes
    for shape in shapes:
        if shape.name == 'month_range_title':
            shape.text = 'Performance Management Report – ' \
                        + calendar.month_name[lastMonth.month] + ' \nHR Operations Quality and Continuous Improvement' 


    # slide 2 (Summary HRC + FO)
    feedbackreceived, rating12Nom, createdC, createdU, resolvedC, resolvedU, \
    openworkload, ageing115, ageing1530, ageing3060, ageing60180, \
    ageingover180, otdNom, otdDenom, QCDoneNom, QCnaNom, QCpassedNom, \
    QCfailedNom, QCDoneNomVol, QCnaNomVol, QCpassedNomVol, QCfailedNomVol = get_summary_data(cursor, sqldate, '%')

    shapes = slides[1].shapes
    for shape in shapes:
        if shape.name == 'month_range':
            text = calendar.month_name[lastMonth.month] + ' ' + str(lastMonth.year)
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'CSATchart':
            chart = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = ['Survey']
            chart_data.add_series('No response',(resolvedC - feedbackreceived ,))
            chart_data.add_series('Satisfied',(rating12Nom,))
            chart_data.add_series('Dissatisfied',(feedbackreceived - rating12Nom,))
            chart.replace_data(chart_data)
            
    #         text = f"""{round((rating12Nom / feedbackreceived) * 100)}% ({rating12Nom}) Very good / Good;
    # {round(100 - (rating12Nom / feedbackreceived) * 100)}% ({feedbackreceived - rating12Nom}) not satisfied
    # Response rate: {round((feedbackreceived * 100) / createdC)}% ({feedbackreceived} answers)"""
    #         text_frame = shape.text_frame
    #         text_frame.clear()
    #         p = text_frame.paragraphs[0]
    #         run = p.add_run()
    #         run.text = text

    #         font = run.font
    #         font.name = 'Arial'
    #         font.size = Pt(14)
    #         font.bold = False
    #         font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'VolumeArea':
            text = f"""{createdU} created / {resolvedU} closed
    """
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'CAArea':
            text = f"""Open workload {openworkload}
    {round(ageing115 * 100 / openworkload)}% ageing <15d / {round(ageing1530 * 100 / openworkload)}% 15-30d
    / {round(ageing3060 * 100 / openworkload)}% 30-60d / {round(ageing60180 * 100 / openworkload)}% 60-180d / {round(ageingover180 * 100 / openworkload)}% >180d"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'OTDArea':
            text = f"""{round(otdNom * 100 / otdDenom)}% On Time Delivery
    """
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'QCAreaTOP':
            text = f"""Out of all survey responses ({feedbackreceived})
    X% (Y) was unsatisfied with Quality"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'QCArea':
            text = f"""{round(QCDoneNom * 100 / (QCDoneNom + QCnaNom))}% ({QCDoneNomVol}) QC done, out of it:
    {round(QCpassedNom * 100 / QCDoneNom)}% ({QCpassedNomVol}) passed / {round(QCfailedNom * 100 / (QCDoneNom + QCnaNom))}% ({QCfailedNomVol}) failed
    {round(QCnaNom * 100 / (QCDoneNom + QCnaNom))} % ({QCnaNomVol}) QC n/a"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

    # slide 5 (Summary HRC)
    feedbackreceived, rating12Nom, createdC, createdU, resolvedC, resolvedU, \
    openworkload, ageing115, ageing1530, ageing3060, ageing60180, \
    ageingover180, otdNom, otdDenom, QCDoneNom, QCnaNom, QCpassedNom, \
    QCfailedNom, QCDoneNomVol, QCnaNomVol, QCpassedNomVol, QCfailedNomVol = get_summary_data(cursor, sqldate, 'Hub_Office')

    shapes = slides[4].shapes
    for shape in shapes:
        if shape.name == 'month_range':
            text = calendar.month_name[lastMonth.month] + ' ' + str(lastMonth.year)
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'CFArea':
            text = f"""{round((rating12Nom / feedbackreceived) * 100)}% ({rating12Nom}) Very good / Good;
    {round(100 - (rating12Nom / feedbackreceived) * 100)}% ({feedbackreceived - rating12Nom}) not satisfied
    Response rate: {round((feedbackreceived * 100) / createdC)}% ({feedbackreceived} answers)"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'VolumeArea':
            text = f"""{createdU} created / {resolvedU} closed
    """
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'CAArea':
            text = f"""Open workload {openworkload}
    {round(ageing115 * 100 / openworkload)}% ageing <15d / {round(ageing1530 * 100 / openworkload)}% 15-30d
    / {round(ageing3060 * 100 / openworkload)}% 30-60d / {round(ageing60180 * 100 / openworkload)}% 60-180d / {round(ageingover180 * 100 / openworkload)}% >180d"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'OTDArea':
            text = f"""{round(otdNom * 100 / otdDenom)}% On Time Delivery
    """
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'QCAreaTOP':
            text = f"""Out of all survey responses ({feedbackreceived})
    X% (Y) was unsatisfied with Quality"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'QCArea':
            text = f"""{round(QCDoneNom * 100 / (QCDoneNom + QCnaNom))}% ({QCDoneNomVol}) QC done, out of it:
    {round(QCpassedNom * 100 / QCDoneNom)}% ({QCpassedNomVol}) passed / {round(QCfailedNom * 100 / (QCDoneNom + QCnaNom))}% ({QCfailedNomVol}) failed
    {round(QCnaNom * 100 / (QCDoneNom + QCnaNom))} % ({QCnaNomVol}) QC n/a"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

    # slide 6 Summary FO
    feedbackreceived, rating12Nom, createdC, createdU, resolvedC, resolvedU, \
    openworkload, ageing115, ageing1530, ageing3060, ageing60180, \
    ageingover180, otdNom, otdDenom, QCDoneNom, QCnaNom, QCpassedNom, \
    QCfailedNom, QCDoneNomVol, QCnaNomVol, QCpassedNomVol, QCfailedNomVol = get_summary_data(cursor, sqldate, 'Front_Office')

    shapes = slides[5].shapes
    names = []
    for shape in shapes:
        if shape.name == 'month_range':
            text = calendar.month_name[lastMonth.month] + ' ' + str(lastMonth.year)
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'CFArea':
            text = f"""{round((rating12Nom / feedbackreceived) * 100)}% ({rating12Nom}) Very good / Good;
    {round(100 - (rating12Nom / feedbackreceived) * 100)}% ({feedbackreceived - rating12Nom}) not satisfied
    Response rate: {round((feedbackreceived * 100) / createdC)}% ({feedbackreceived} answers)"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'VolumeArea':
            text = f"""{createdU} created / {resolvedU} closed
    """
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'CAArea':
            text = f"""Open workload {openworkload}
    {round(ageing115 * 100 / openworkload)}% ageing <15d / {round(ageing1530 * 100 / openworkload)}% 15-30d
    / {round(ageing3060 * 100 / openworkload)}% 30-60d / {round(ageing60180 * 100 / openworkload)}% 60-180d / {round(ageingover180 * 100 / openworkload)}% >180d"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'OTDArea':
            text = f"""{round(otdNom * 100 / otdDenom)}% On Time Delivery
    """
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'QCAreaTOP':
            text = f"""Out of all survey responses ({feedbackreceived})
    X% (Y) was unsatisfied with Quality"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)

        elif shape.name == 'QCArea':
            text = f"""{round(QCDoneNom * 100 / (QCDoneNom + QCnaNom))}% ({QCDoneNomVol}) QC done, out of it:
    {round(QCpassedNom * 100 / QCDoneNom)}% ({QCpassedNomVol}) passed / {round(QCfailedNom * 100 / (QCDoneNom + QCnaNom))}% ({QCfailedNomVol}) failed 
    {round(QCnaNom * 100 / (QCDoneNom + QCnaNom))} % ({QCnaNomVol}) QC n/a"""
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text

            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)
            
    print('Summary slides ready')


    # slides 11-20 (CSAT)
    for i in range(10, 20):
        shapes = slides[i].shapes
        shapes = [x for x in shapes if x.name[-1] == 'a']
        
        for shape in shapes:
            if shape.name[0:3] in ('TA_', 'TM_', 'GM_', 'PA_', 'ELC', 'CB_',
                                'PY_', 'LD_', 'BGL', 'KRK', 'SPL', 'TLL',
                                'CSE', 'NEU', 'NAM', 'SAM', 'NEA', 'MEA', 'SEA'):
                months, VeryGood, Good, Dissatisfied, HighlyDissatisfied, ResponseRate = get_cs_chart(cursor, sqldate, shape.name[0:3])
                
                
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('Very Good', VeryGood)
                chart_data.add_series('Good', Good)
                chart_data.add_series('Dissatisfied', Dissatisfied)
                chart_data.add_series('Highly Dissatisfied', HighlyDissatisfied)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[i]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_STACKED_100, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(8)
                chart.chart_style = 1
                
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(255, 255, 255)
                data_labels.position = XL_LABEL_POSITION.CENTER
            

                value_axis = chart.value_axis
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(10)
                
            
                
    print('CSAT slides ready')

    # slides 21-25 (OTD)
    for i in range(20, 25):
        shapes = slides[i].shapes
        shapes = [x for x in shapes if x.name[-1] == 'a']
        for shape in shapes:
            if shape.name[0:3] in ('TA_', 'TM_', 'GM_', 'PA_', 'ELC', 'CB_',
                                'PY_', 'LD_', 'BGL', 'KRK', 'SPL', 'TLL',
                                'CSE', 'NEU', 'NAM', 'SAM', 'NEA', 'MEA', 'SEA'):
                months, data = get_otd_chart(cursor, sqldate, shape.name[0:3])

                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('Series 1', data)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[i]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0"%"'

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
                
    print('OTD slides ready')


    # slides 26-30 (Quality)
    for i in range(25, 30):
        shapes = slides[i].shapes
        shapes = [x for x in shapes if x.name[-1] == 'a']
        for shape in shapes:
            if shape.name[0:3] in ('TA_', 'TM_', 'GM_', 'PA_', 'ELC', 'CB_',
                                'PY_', 'LD_', 'BGL', 'KRK', 'SPL', 'TLL',
                                'CSE', 'NEU', 'NAM', 'SAM', 'NEA', 'MEA', 'SEA'):
                months, data = get_qc_chart(cursor, sqldate, shape.name[0:3])

                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('Series 1', data)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[i]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(0, 0, 0)
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                data_labels.number_format = '0.0"%"'

                value_axis = chart.value_axis
                value_axis.maximum_scale = 100.0
                value_axis.minimum_scale = 0.0
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
                
    print('Quality slides ready')


    # slides 31-35 (Workload)
    for i in range(30, 35):
        shapes = slides[i].shapes
        shapes = [x for x in shapes if x.name[-1] == 'a']
        for shape in shapes:
            if shape.name[0:3] in ('TA_', 'TM_', 'GM_', 'PA_', 'ELC', 'CB_',
                                'PY_', 'LD_', 'BGL', 'KRK', 'SPL', 'TLL',
                                'CSE', 'NEU', 'NAM', 'SAM', 'NEA', 'MEA', 'SEA'):
                months, created, resolved = get_vl_chart(cursor, sqldate, shape.name[0:3])

                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('Created', created)
                chart_data.add_series('Resolved', resolved)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[i]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(8)
                #chart.chart_style = 1

                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(255, 255, 255)
                data_labels.position = XL_LABEL_POSITION.CENTER

                value_axis = chart.value_axis
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)
                value_axis.has_major_gridlines = False
                chart.value_axis.visible = False

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
                
    print('Workload slides ready')


    # slides 36-45 (Case Ageing)
    for i in range(35, 45):
        shapes = slides[i].shapes
        shapes = [x for x in shapes if x.name[-1] == 'a']
        for shape in shapes:
            if shape.name[0:3] in ('TA_', 'TM_', 'GM_', 'PA_', 'ELC', 'CB_', 'PY_', 'LD_', 'BGL', 'KRK', 'SPL', 'TLL',
                                    'CSE', 'NEU', 'NAM', 'SAM', 'NEA', 'MEA', 'SEA'):
                months, active115, active1530, active3060, activeover60, longrunning = get_ca_chart(cursor, sqldate, shape.name[0:3])
                
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
            
                for c in chart_data.categories:
                    c.add_sub_category('active')
                    c.add_sub_category('on hold')
                

                
                chart_data.add_series('Long running cases', longrunning)
                chart_data.add_series('1 - 15 days', active115)
                chart_data.add_series('15 - 30 days', active1530)
                chart_data.add_series('30 - 60 days', active3060)
                chart_data.add_series('over 60 days', activeover60)
                
            
                
                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[i]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
                )
                
                # edit look of the slide
                chart = graphic_frame.chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(8)
                # chart.chart_style = 1
                

                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(255, 255, 255)
                data_labels.position = XL_LABEL_POSITION.CENTER

                value_axis = chart.value_axis
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(8)
                
    print('Case Ageing slides ready')

                

    # slides 46-48 (Case Ageing Long Term)
    for i in range(45, 48):
        shapes = slides[i].shapes
        shapes = [x for x in shapes if x.name[-1] == 'a']
        for shape in shapes:
            if shape.name[0:3] in ('TA_', 'TM_', 'GM_', 'PA_', 'ELC', 'CB_', 'PY_', 'LD_', 'BGL', 'KRK', 'SPL', 'TLL',
                                    'CSE', 'NEU', 'NAM', 'SAM', 'NEA', 'MEA', 'SEA'):
                months, ag120150, ag150180, agover180  = get_calt_chart(cursor, sqldate, shape.name[0:3])
                
                # define chart data
                chart_data = ChartData()
                chart_data.categories = [months[0], months[1], months[2],
                                        months[3], months[4], months[5]]
                chart_data.add_series('120 - 150 days', ag120150)
                chart_data.add_series('150 - 180 days', ag150180)
                chart_data.add_series('over 180 days', agover180)

                # add chart to slide
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                slide = slides[i]
                graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
                )

                # edit look of the slide
                chart = graphic_frame.chart
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(8)
                # chart.chart_style = 1

                data_labels.font.size = Pt(6)
                data_labels.font.color.rgb = RGBColor(255, 255, 255)
                # data_labels.position = XL_LABEL_POSITION.INSIDE_END
                data_labels.number_format = '0"%"'

                value_axis = chart.value_axis
                # value_axis.maximum_scale = 100.0
                # value_axis.minimum_scale = 0.0
                chart.value_axis.visible = False
                value_axis.has_major_gridlines = False
                tick_labels = value_axis.tick_labels
                tick_labels.font.size = Pt(10)

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(7)
                
    print('Case Ageing LT slides ready')


    # save presentation
    name = 'HR Operations Monthly Report ' + calendar.month_name[lastMonth.month] + ' ' + str(
        lastMonth.year) + '.pptx'
    
    if not os.path.isdir(f"{os.getcwd()}\\Created\\"):
        os.mkdir(f"{os.getcwd()}\\Created\\")

    prs.save(f"{os.getcwd()}\\Created\\" + name)
    print(name, 'created')
    print("--- %s seconds ---" % (time.time() - start_time))
    cursor.close()
    conn.close()