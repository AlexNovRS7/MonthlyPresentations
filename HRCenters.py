from pptx import Presentation
import calendar
import os 

from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION

from pptx.dml.color import RGBColor
from pptx.util import Pt

from countries import Centers, units

import HRCenterscharts

import datetime
import time

import pyodbc



today = datetime.date.today()
first = today.replace(day=1)
sqldate = f'{first.year}-{first.month - 1}-{first.day}' # date for SQL querries
lastMonth = first - datetime.timedelta(days=1)



def CreateReport(HRCenter):
    t1 = time.perf_counter()

    conn = pyodbc.connect('Driver={SQL Server};'
                      'Server=XE-S-CWPDB01P.XE.ABB.COM;'
                      'Database=PG_AskHR_KPI_2;'
                      'Trusted_Connection=yes;')

    cursor = conn.cursor()

    global HRCenters
    print("Creating: ", Centers[HRCenter])
    


    prs = Presentation(f"{os.getcwd()}\\templates\\Monthly_HR_Center_Report_TEMPLATE_2.0.pptx")
    slides = prs.slides

    # slide 1 (Cover)
    slide = slides[0]
    shapes = slide.shapes
    for shape in shapes:
        if shape.name == 'ctry_range_title':
            # shape.text = 'Performance Management Report – ' +  Centers[HRCenter] + '\nHR Operations Quality and Continuous Improvement'
            
            text = 'Performance Management Report – ' +  Centers[HRCenter] + '\nHR Operations Quality and Continuous Improvement'
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text[32:32+len(Centers[HRCenter])]
            
            font = run.font
            font.name = 'Arial'
            font.size = Pt(24)
            font.bold = True
            font.color.rgb = RGBColor(255, 0, 0)
        
        if shape.name == 'date_range':
            text = str(today)
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            
            font = run.font
            font.name = 'Arial'
            font.size = Pt(14)
            font.bold = False
            font.color.rgb = RGBColor(0, 0, 0)
    
    # slide 2 (Exec Summary)

    # save presentation
    name = 'GBS HR Monthly Review ' + Centers[HRCenter] + ' ' + calendar.month_name[lastMonth.month] + ' ' + str(
        lastMonth.year) + '.pptx'

    if not os.path.isdir(f"{os.getcwd()}\\Created\\"):
        os.mkdir(f"{os.getcwd()}\\Created\\")
        os.mkdir(f"{os.getcwd()}\\Created\\HRCenters\\")

    if not os.path.isdir(f"{os.getcwd()}\\Created\\HRCenters\\"):
        os.mkdir(f"{os.getcwd()}\\Created\\HRCenters\\")
        
    prs.save(f"{os.getcwd()}\\Created\\HRCenters\\" + name)
        
    t2 = time.perf_counter()
    print(Centers[HRCenter], f' created. \t time: {round((t2 - t1), 2)} seconds')

    cursor.close()